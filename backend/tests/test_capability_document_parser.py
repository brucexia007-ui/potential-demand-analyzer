"""能力资料解析保留结构来源并生成稳定切片。"""
from __future__ import annotations

import subprocess
from io import BytesIO

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter
from weasyprint import HTML

from app.capabilities import document_parser
from app.capabilities.document_parser import (
    CapabilityDocumentParseError,
    MAX_PDF_BYTES,
    ParsedSegment,
    chunk_segments,
    parse_capability_document,
)


def test_parse_docx_preserves_heading() -> None:
    document = Document()
    document.add_heading("智能客服平台", level=1)
    document.add_paragraph("支持多渠道统一接入。")
    buffer = BytesIO()
    document.save(buffer)

    segments = parse_capability_document(
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        content=buffer.getvalue(),
    )

    assert segments == [ParsedSegment(content="支持多渠道统一接入。", heading="智能客服平台")]


def test_parse_xlsx_preserves_sheet_reference() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "产品能力"
    sheet.append(["能力", "限制"])
    sheet.append(["智能质检", "仅支持中文"])
    buffer = BytesIO()
    workbook.save(buffer)

    segments = parse_capability_document(
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        content=buffer.getvalue(),
    )

    assert segments[0].page_ref == "sheet:产品能力"
    assert "智能质检 | 仅支持中文" in segments[0].content


def test_text_decode_and_deterministic_chunking() -> None:
    parsed = parse_capability_document(mime_type="text/plain", content="第一段\n第二段".encode("gb18030"))
    chunks = chunk_segments([ParsedSegment("a" * 250 + "\n" + "b" * 250, page_ref="page:1")], max_chars=200)

    assert parsed[0].content == "第一段\n第二段"
    assert [len(item.content) for item in chunks] == [200, 50, 200, 50]
    assert all(item.page_ref == "page:1" for item in chunks)


def test_parse_pptx_preserves_slide_reference() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Customer service solution"
    slide.placeholders[1].text = "Supports omnichannel access"
    buffer = BytesIO()
    presentation.save(buffer)

    segments = parse_capability_document(
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        content=buffer.getvalue(),
    )

    assert segments[0].page_ref == "slide:1"
    assert segments[0].heading == "Customer service solution"
    assert "Supports omnichannel access" in segments[0].content


def test_parse_pdf_preserves_page_reference() -> None:
    content = HTML(string="<p>Product capability statement</p>").write_pdf()

    segments = parse_capability_document(mime_type="application/pdf", content=content)

    assert segments[0].page_ref == "page:1"
    assert "Product capability statement" in segments[0].content


def test_parse_pdf_rejects_payload_over_byte_limit(monkeypatch: pytest.MonkeyPatch) -> None:
    def unexpected_worker_call(_: bytes) -> list[ParsedSegment]:
        raise AssertionError("超限文件不应进入 PDF 解析子进程")

    monkeypatch.setattr(document_parser, "_run_pdf_worker", unexpected_worker_call)

    with pytest.raises(CapabilityDocumentParseError, match="大小超过"):
        parse_capability_document(mime_type="application/pdf", content=b"x" * (MAX_PDF_BYTES + 1))


def test_parse_pdf_reports_worker_timeout(monkeypatch: pytest.MonkeyPatch) -> None:
    def timeout(*args: object, **kwargs: object) -> subprocess.CompletedProcess[bytes]:
        raise subprocess.TimeoutExpired(cmd=["python", "worker"], timeout=20)

    monkeypatch.setattr(document_parser.subprocess, "run", timeout)

    with pytest.raises(CapabilityDocumentParseError, match="解析超时"):
        parse_capability_document(mime_type="application/pdf", content=b"%PDF-test")


def test_parse_pdf_rejects_excessive_page_count() -> None:
    writer = PdfWriter()
    for _ in range(document_parser.MAX_PDF_PAGES + 1):
        writer.add_blank_page(width=72, height=72)
    buffer = BytesIO()
    writer.write(buffer)

    with pytest.raises(CapabilityDocumentParseError, match="页数超过"):
        parse_capability_document(mime_type="application/pdf", content=buffer.getvalue())
