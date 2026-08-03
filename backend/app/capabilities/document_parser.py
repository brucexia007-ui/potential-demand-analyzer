"""企业能力资料的结构化文本解析与确定性切片。"""
from __future__ import annotations

import csv
import json
import os
import subprocess
import sys
from dataclasses import dataclass
from io import BytesIO, StringIO


MAX_PDF_BYTES = 20 * 1024 * 1024
MAX_PDF_PAGES = 500
MAX_PDF_EXTRACTED_CHARS = 5_000_000
PDF_PARSE_TIMEOUT_SECONDS = 20


class CapabilityDocumentParseError(ValueError):
    pass


@dataclass(frozen=True)
class ParsedSegment:
    content: str
    page_ref: str | None = None
    heading: str | None = None


def parse_capability_document(*, mime_type: str, content: bytes) -> list[ParsedSegment]:
    try:
        if mime_type == "application/pdf":
            return _parse_pdf(content)
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _parse_docx(content)
        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _parse_pptx(content)
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return _parse_xlsx(content)
        if mime_type in {"text/plain", "text/markdown"}:
            text = _decode_text(content)
            return [ParsedSegment(content=text)] if text.strip() else []
        if mime_type == "text/csv":
            return _parse_csv(content)
    except CapabilityDocumentParseError:
        raise
    except Exception as error:
        raise CapabilityDocumentParseError("能力资料解析失败") from error
    raise CapabilityDocumentParseError("不支持的能力资料类型")


def chunk_segments(segments: list[ParsedSegment], *, max_chars: int = 1800) -> list[ParsedSegment]:
    if max_chars < 200:
        raise ValueError("切片长度不能小于 200 字符")
    chunks: list[ParsedSegment] = []
    for segment in segments:
        paragraphs = [item.strip() for item in segment.content.replace("\r\n", "\n").split("\n") if item.strip()]
        current = ""
        for paragraph in paragraphs:
            if len(paragraph) > max_chars:
                if current:
                    chunks.append(ParsedSegment(current, segment.page_ref, segment.heading))
                    current = ""
                for start in range(0, len(paragraph), max_chars):
                    chunks.append(ParsedSegment(paragraph[start:start + max_chars], segment.page_ref, segment.heading))
                continue
            candidate = f"{current}\n{paragraph}".strip()
            if current and len(candidate) > max_chars:
                chunks.append(ParsedSegment(current, segment.page_ref, segment.heading))
                current = paragraph
            else:
                current = candidate
        if current:
            chunks.append(ParsedSegment(current, segment.page_ref, segment.heading))
    return chunks


def _parse_pdf(content: bytes) -> list[ParsedSegment]:
    if len(content) > MAX_PDF_BYTES:
        raise CapabilityDocumentParseError("PDF 文件大小超过 20 MB 处理上限")
    return _run_pdf_worker(content)


def _run_pdf_worker(content: bytes) -> list[ParsedSegment]:
    worker_env = {
        "PYTHONPATH": os.pathsep.join(str(path) for path in sys.path if path),
        "PYTHONUNBUFFERED": "1",
        "PDF_PARSE_MAX_BYTES": str(MAX_PDF_BYTES),
        "PDF_PARSE_MAX_PAGES": str(MAX_PDF_PAGES),
        "PDF_PARSE_MAX_CHARS": str(MAX_PDF_EXTRACTED_CHARS),
    }
    try:
        result = subprocess.run(
            [sys.executable, "-m", "app.capabilities.pdf_parser_worker"],
            input=content,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=False,
            timeout=PDF_PARSE_TIMEOUT_SECONDS,
            env=worker_env,
            start_new_session=True,
        )
    except subprocess.TimeoutExpired as error:
        raise CapabilityDocumentParseError("PDF 解析超时") from error

    if result.returncode != 0:
        reason = result.stderr.decode("utf-8", errors="replace").strip()
        if reason == "PDF_PAGE_LIMIT":
            raise CapabilityDocumentParseError(f"PDF 页数超过 {MAX_PDF_PAGES} 页处理上限")
        if reason == "PDF_TEXT_LIMIT":
            raise CapabilityDocumentParseError("PDF 提取文本超过处理上限")
        if reason == "PDF_SIZE_LIMIT":
            raise CapabilityDocumentParseError("PDF 文件大小超过 20 MB 处理上限")
        if reason == "PDF_MEMORY_LIMIT":
            raise CapabilityDocumentParseError("PDF 解析内存超过处理上限")
        raise CapabilityDocumentParseError("PDF 文件无法安全解析")

    try:
        payload = json.loads(result.stdout.decode("utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError) as error:
        raise CapabilityDocumentParseError("PDF 解析结果无效") from error
    if not isinstance(payload, list):
        raise CapabilityDocumentParseError("PDF 解析结果无效")

    segments: list[ParsedSegment] = []
    for item in payload:
        if not isinstance(item, dict):
            raise CapabilityDocumentParseError("PDF 解析结果无效")
        text = item.get("content")
        page_ref = item.get("page_ref")
        if not isinstance(text, str) or not isinstance(page_ref, str):
            raise CapabilityDocumentParseError("PDF 解析结果无效")
        segments.append(ParsedSegment(content=text, page_ref=page_ref))
    return segments


def _parse_docx(content: bytes) -> list[ParsedSegment]:
    from docx import Document

    document = Document(BytesIO(content))
    segments: list[ParsedSegment] = []
    heading: str | None = None
    buffer: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style and paragraph.style.name.lower().startswith("heading"):
            if buffer:
                segments.append(ParsedSegment("\n".join(buffer), heading=heading))
                buffer = []
            heading = text
        else:
            buffer.append(text)
    if buffer:
        segments.append(ParsedSegment("\n".join(buffer), heading=heading))
    return segments


def _parse_pptx(content: bytes) -> list[ParsedSegment]:
    from pptx import Presentation

    presentation = Presentation(BytesIO(content))
    result: list[ParsedSegment] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if texts:
            result.append(ParsedSegment("\n".join(texts), page_ref=f"slide:{index}", heading=texts[0]))
    return result


def _parse_xlsx(content: bytes) -> list[ParsedSegment]:
    from openpyxl import load_workbook

    workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
    result: list[ParsedSegment] = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            result.append(ParsedSegment("\n".join(rows), page_ref=f"sheet:{sheet.title}", heading=sheet.title))
    workbook.close()
    return result


def _parse_csv(content: bytes) -> list[ParsedSegment]:
    rows = [" | ".join(cell.strip() for cell in row) for row in csv.reader(StringIO(_decode_text(content)))]
    text = "\n".join(row for row in rows if row.strip(" |"))
    return [ParsedSegment(content=text)] if text else []


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise CapabilityDocumentParseError("文本编码无法识别")
